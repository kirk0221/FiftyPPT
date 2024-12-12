# 모듈 import
import os
import cv2
import numpy as np
import requests
import json
import uuid
import time
from ultralytics import YOLO
from pptx import Presentation
from pptx.util import Inches, Pt
import random
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from createppt.forms import PPTForm
from django.contrib.auth.decorators import login_required
from django.shortcuts import get_object_or_404
from .models import PPT, Cart
from django.http import FileResponse, Http404, JsonResponse
from django.shortcuts import render, redirect


def process_image_and_generate_ppt(secret_key, api_url, image_file, model_path, output_dir, ppt_output_path):
    response_data, img_dimensions = ocr_and_save_data(secret_key, api_url, image_file)
    if response_data is None:
        print(f"{image_file}의 OCR 처리 중 오류가 발생했습니다.")
        return

    mask = create_mask_from_ocr(response_data, img_dimensions)
    inpainted_image_path = inpaint_image(image_file, mask, output_dir)

    # 객체 검출 모델 로드 및 예측
    model = YOLO(model_path)
    detection_results = model.predict(
        source=inpainted_image_path,
        imgsz=640,
        conf=0.5,
        iou=0.4,
        save=False
    )

    save_presentation_with_text_and_shapes(response_data, img_dimensions, ppt_output_path, detection_results)


def merge_close_texts(response_data, max_distance=50):
    """
    OCR 결과에서 가까운 텍스트를 병합합니다.
    - max_distance: 텍스트 간 최대 거리 (픽셀 단위)
    """
    for image_info in response_data['images']:
        fields = image_info['fields']
        merged_fields = []

        while fields:
            base_field = fields.pop(0)
            base_vertices = base_field['boundingPoly']['vertices']
            base_text = base_field['inferText']

            # 병합 가능한 텍스트 찾기
            for other_field in fields[:]:
                other_vertices = other_field['boundingPoly']['vertices']
                distance = abs(base_vertices[2]['x'] - other_vertices[0]['x'])

                if distance <= max_distance:  # 병합 기준
                    base_text += " " + other_field['inferText']
                    # 병합된 영역 업데이트
                    base_vertices[2]['x'] = max(base_vertices[2]['x'], other_vertices[2]['x'])
                    base_vertices[2]['y'] = max(base_vertices[2]['y'], other_vertices[2]['y'])
                    fields.remove(other_field)

            # 병합 결과 추가
            merged_fields.append({
                'inferText': base_text,
                'boundingPoly': {'vertices': base_vertices}
            })

        # 업데이트된 필드로 교체
        image_info['fields'] = merged_fields

    return response_data


def ocr_and_save_data(secret_key, api_url, image_file):
    """
    이미지에서 OCR을 사용하여 텍스트 위치 및 폰트 크기 정보를 추출하고 결과 데이터를 반환합니다.
    """
    image = cv2.imread(image_file, cv2.IMREAD_GRAYSCALE)
    if image is None:
        print("이미지를 불러올 수 없습니다.")
        return None, None

    img_height, img_width = image.shape

    # CLOVA OCR API 요청 JSON
    request_json = {
        'images': [{'format': 'jpg', 'name': 'demo'}],
        'requestId': str(uuid.uuid4()),
        'version': 'V2',
        'timestamp': int(round(time.time() * 1000)),
        'options': {
            'lineMerge': True,
            'paragraphMerge': True,
            'language': 'ko',
            'detectOrientation': True,

        }
    }

    # API 요청
    payload = {'message': json.dumps(request_json).encode('UTF-8')}
    files = [('file', open(image_file, 'rb'))]
    headers = {'X-OCR-SECRET': secret_key}

    response = requests.post(api_url, headers=headers, data=payload, files=files)
    response_data = json.loads(response.text)

    # 병합 후처리 적용
    response_data = merge_close_texts(response_data, max_distance=10)

    return response_data, (img_height, img_width)


def create_mask_from_ocr(response_data, img_dimensions):
    """
    OCR 결과를 사용하여 마스크를 생성합니다.
    """
    img_height, img_width = img_dimensions
    mask = np.zeros((img_height, img_width), dtype=np.uint8)
    for image_info in response_data['images']:
        for field in image_info['fields']:
            vertices = field['boundingPoly']['vertices']
            x1, y1 = int(vertices[0]['x']), int(vertices[0]['y'])
            x2, y2 = int(vertices[2]['x']), int(vertices[2]['y'])
            cv2.rectangle(mask, (x1, y1), (x2, y2), (255), thickness=cv2.FILLED)
    return mask


def inpaint_image(image_file, mask, output_dir):
    """
    마스크를 사용하여 인페인팅된 이미지를 저장합니다.
    """
    image_color = cv2.imread(image_file)
    if image_color is None:
        print("컬러 이미지를 불러올 수 없습니다.")
        return None

    inpainted_image = cv2.inpaint(image_color, mask, 7, cv2.INPAINT_TELEA)

    inpaint_output_path = os.path.join(output_dir, 'inpainted_image.jpg')
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    cv2.imwrite(inpaint_output_path, inpainted_image)
    print(f"인페인팅된 이미지가 {inpaint_output_path}에 성공적으로 저장되었습니다.")
    return inpaint_output_path


def save_presentation_with_text_and_shapes(response_data, img_dimensions, ppt_output_path, detection_results):
    """OCR 결과와 객체 검출 결과를 바탕으로 PPT 저장"""
    img_height, img_width = img_dimensions
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # 빈 슬라이드 레이아웃 선택
    slide = prs.slides.add_slide(slide_layout)

    # PPT 슬라이드 크기
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches

    # 이미지와 슬라이드 비율 계산
    x_scale = slide_width / img_width
    y_scale = slide_height / img_height
    scale = min(x_scale, y_scale)

    # 도형을 크기 순으로 정렬
    sorted_boxes = sorted(detection_results[0].boxes,
                          key=lambda box: (box.xyxy[0][2] - box.xyxy[0][0]) * (box.xyxy[0][3] - box.xyxy[0][1]),
                          reverse=True)

    # 정렬된 도형을 추가
    for pred in sorted_boxes:
        x1, y1, x2, y2 = map(int, pred.xyxy[0].tolist())
        class_id = int(pred.cls[0])
        left = Inches(x1 * scale)
        top = Inches(y1 * scale)
        width_in_inches = Inches((x2 - x1) * scale)
        height_in_inches = Inches((y2 - y1) * scale)

        # 기본 도형 색상
        fill_color = RGBColor(255, 255, 255)  # 기본 흰색
        skip_text = False  # 텍스트 상자 추가 여부 플래그

        # 도형 내부 OCR 텍스트 확인
        for image_info in response_data['images']:
            for field in image_info['fields']:
                vertices = field['boundingPoly']['vertices']
                text_x1, text_y1 = int(vertices[0]['x']), int(vertices[0]['y'])
                text_x2, text_y2 = int(vertices[2]['x']), int(vertices[2]['y'])

                # OCR 텍스트가 도형 내부에 있는지 확인
                if x1 <= text_x1 <= x2 and y1 <= text_y1 <= y2:
                    text = field['inferText']
                    if "연한파랑" in text or "파랑" in text:
                        fill_color = RGBColor(173, 216, 230)  # 연한 파랑
                        skip_text = True
                    elif "연한빨강" in text or "빨강" in text:
                        fill_color = RGBColor(255, 182, 193)  # 연한 빨강
                        skip_text = True
                    elif "연한초록" in text or "초록" in text:
                        fill_color = RGBColor(144, 238, 144)  # 연한 초록
                        skip_text = True
                    elif "연한노랑" in text or "노랑" in text:
                        fill_color = RGBColor(255, 255, 224)  # 연한 노랑
                        skip_text = True
                    elif "연한주황" in text or "주황" in text:
                        fill_color = RGBColor(255, 200, 124)  # 연한 초록
                        skip_text = True
                    elif "연한보라" in text or "보라" in text:
                        fill_color = RGBColor(216, 191, 216)  # 연한 노랑
                        skip_text = True

        if class_id == 0:  # 사각형
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width_in_inches, height_in_inches)
        elif class_id == 1:  # 원
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width_in_inches, height_in_inches)
        elif class_id == 2:  # 둥근 사각형
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width_in_inches, height_in_inches)
        elif class_id == 3:  # 마름모
            shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, left, top, width_in_inches, height_in_inches)
        elif class_id == 4:  # 위쪽 삼각형
            shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width_in_inches, height_in_inches)
        elif class_id == 5:  # 오른쪽 화살표
            shape = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, left, top, width_in_inches, height_in_inches)
            shape.rotation = 90

        # 도형 스타일 설정
        shape.line.color.rgb = RGBColor(0, 0, 0)  # 검은색 테두리
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color  # OCR 텍스트에 따른 색상

    # OCR 결과를 텍스트 상자로 추가
    for image_info in response_data['images']:
        for field in image_info['fields']:
            text = field['inferText']
            vertices = field['boundingPoly']['vertices']
            x1, y1 = int(vertices[0]['x']), int(vertices[0]['y'])
            x2, y2 = int(vertices[2]['x']), int(vertices[2]['y'])

            left = Inches(x1 * scale)
            top = Inches(y1 * scale)
            width = Inches((x2 - x1) * scale)
            height = Inches((y2 - y1) * scale)

            # 텍스트 상자를 특정 색상 텍스트가 아닌 경우에만 추가
            if not any(color in text for color in ["파랑", "빨강", "초록", "노랑", "보라", "주황"]):
                # 텍스트 상자 추가
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.text = text

                # 폰트 크기를 고정된 15 포인트로 설정
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(15)

    prs.save(ppt_output_path)
    print(f"PPT가 {ppt_output_path}에 성공적으로 저장되었습니다.")


def create_ppt(request):
    if request.method == 'POST':
        form = PPTForm(request.POST, request.FILES)
        if form.is_valid():
            title = form.cleaned_data['title']
            image_file = form.cleaned_data['image']  # 업로드된 이미지 파일

            image_path = os.path.join('image_file', image_file.name)
            with open(os.path.join('media', image_path), 'wb') as f:
                for chunk in image_file.chunks():
                    f.write(chunk)

            ppt_output_path = os.path.join('ppt_output', f"{title}.pptx")

            # PPT 생성 함수 실행
            process_image_and_generate_ppt(
                secret_key='',  # 네이버 OCR 시크릿 키  ***********************************************************************
                api_url='',  # 네이버 OCR API url   ***********************************************************************
                model_path='/home/fiftyppt/fiftyPPT_project/model/best.pt',
                image_file=os.path.join('media', image_path),  # 저장된 이미지 경로
                output_dir='/home/fiftyppt/fiftyPPT_project/media/output',
                ppt_output_path=os.path.join('media', ppt_output_path)  # ppt 출력 경로
            )

            ppt = PPT.objects.create(
                user=request.user,
                title=title,
                file_path=ppt_output_path,  # 'media/'를 제외한 상대 경로만 저장
                image=image_file,
            )
            user_ppts = PPT.objects.filter(user=request.user)
            return render(request, 'Mentor/myppt.html', {'ppts': user_ppts})

    else:
        form = PPTForm()

    return render(request, 'Mentor/createppt.html', {'form': form})


@login_required(login_url='/ppt/login_required/')
def myppt(request):
    user_ppts = PPT.objects.filter(user=request.user)
    return render(request, 'Mentor/myppt.html', {'ppts': user_ppts})


def download_ppt(request, ppt_id):
    ppt = get_object_or_404(PPT, id=ppt_id)

    ppt_file_path = ppt.file_path.path  # 실제 파일 경로를 추출
    print("다운로드하려는 피피티 절대 경로", ppt_file_path)
    if not os.path.exists(ppt_file_path):
        raise Http404("PPT 파일을 찾을 수 없습니다.")

    response = FileResponse(open(ppt_file_path, 'rb'),
                            content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = f'attachment; filename="{os.path.basename(ppt_file_path)}"'
    return response


def all_ppt_list(request):
    allppts = PPT.objects.filter(is_public=True)

    if request.user.is_authenticated:
        cart_ppt_ids = Cart.objects.filter(user=request.user).values_list('ppt_id', flat=True)
    else:
        cart_ppt_ids = []

    # 템플릿에 ppt 목록과 장바구니 상태를 전달합니다.
    return render(request, 'Mentor/allppt.html', {'allppts': allppts, 'cart_ppt_ids': cart_ppt_ids})


def add_to_cart(request, ppt_id):
    user = request.user
    if not user.is_authenticated:
        return JsonResponse({'success': False, 'message': '로그인이 필요합니다.'}, status=400)

    try:
        ppt = PPT.objects.get(id=ppt_id)
        if Cart.objects.filter(user=user, ppt=ppt).exists():
            return JsonResponse({'success': False, 'message': '이미 장바구니에 추가된 PPT입니다.'})

        Cart.objects.create(user=user, ppt=ppt)
        return JsonResponse({'success': True, 'message': '장바구니에 추가되었습니다!'})
    except PPT.DoesNotExist:
        return JsonResponse({'success': False, 'message': 'PPT를 찾을 수 없습니다.'}, status=404)


@login_required(login_url='/ppt/login_required/')
def my_cart(request):
    if not request.user.is_authenticated:
        return render(request, 'Mentor/register.html')

    user = request.user

    ppts = Cart.objects.filter(user=user).values_list('ppt', flat=True)
    ppt_list = PPT.objects.filter(id__in=ppts)

    return render(request, 'Mentor/mycart.html', {'ppts': ppt_list})


@login_required
def delete_ppt(request, ppt_id):
    ppt = get_object_or_404(PPT, id=ppt_id, user=request.user)
    ppt.delete()

    return redirect('/')


def login_required_view(request):
    return render(request, 'Mentor/login_required.html')
